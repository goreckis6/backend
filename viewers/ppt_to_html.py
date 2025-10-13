#!/usr/bin/env python3
"""
PPT/PPTX (PowerPoint Presentation) to HTML converter for web preview.
Uses python-pptx for direct PPTX parsing, LibreOffice as fallback.
Supports both legacy PPT and modern PPTX formats.
"""

import argparse
import os
import sys
import subprocess
import shutil
import traceback

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("WARNING: python-pptx not available")

def convert_pptx_to_html_direct(pptx_file, html_file):
    """
    Convert PPTX to HTML using python-pptx library (direct parsing).
    Only works for PPTX format (not legacy PPT).
    
    Args:
        pptx_file (str): Path to input PPTX file
        html_file (str): Path to output HTML file
    
    Returns:
        bool: True if conversion successful, False otherwise
    """
    if not PPTX_AVAILABLE:
        print("ERROR: python-pptx not available")
        return False
    
    print(f"Attempting PPTX to HTML conversion with python-pptx...")
    
    try:
        # Load presentation
        prs = Presentation(pptx_file)
        
        # Get presentation info
        slide_count = len(prs.slides)
        print(f"Loaded presentation: {slide_count} slides")
        
        # Build HTML
        html_parts = []
        html_parts.append("""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PowerPoint Preview</title>
    <style>
        body { 
            margin: 0; 
            padding: 0; 
            font-family: Arial, sans-serif;
            background: #f5f5f5;
        }
        #header-bar {
            position: sticky;
            top: 0;
            background: linear-gradient(135deg, #0078d4 0%, #00a4ef 100%);
            color: white;
            padding: 12px 20px;
            display: flex;
            justify-content: space-between;
            align-items: center;
            box-shadow: 0 2px 8px rgba(0,0,0,0.1);
            z-index: 1000;
        }
        #header-bar h1 {
            margin: 0;
            font-size: 18px;
            font-weight: 600;
        }
        .header-actions {
            display: flex;
            gap: 10px;
        }
        .header-btn {
            background: rgba(255,255,255,0.2);
            border: 1px solid rgba(255,255,255,0.3);
            color: white;
            padding: 8px 16px;
            border-radius: 6px;
            cursor: pointer;
            font-size: 14px;
            font-weight: 500;
            transition: all 0.2s;
        }
        .header-btn:hover {
            background: rgba(255,255,255,0.3);
            transform: translateY(-1px);
        }
        .container {
            max-width: 1200px;
            margin: 20px auto;
            padding: 20px;
        }
        .slide {
            background: white;
            border-radius: 8px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.1);
            margin-bottom: 30px;
            padding: 30px;
            page-break-after: always;
        }
        .slide-header {
            border-bottom: 2px solid #0078d4;
            padding-bottom: 10px;
            margin-bottom: 20px;
        }
        .slide-number {
            color: #0078d4;
            font-weight: bold;
            font-size: 14px;
        }
        .slide-title {
            font-size: 28px;
            font-weight: bold;
            color: #333;
            margin: 10px 0;
        }
        .slide-content {
            line-height: 1.6;
            color: #555;
        }
        .text-box {
            margin-bottom: 15px;
        }
        .shape-list {
            list-style-position: inside;
            margin: 10px 0;
        }
        @media print {
            #header-bar { display: none; }
            .container { margin: 0; padding: 0; max-width: none; }
            .slide { box-shadow: none; margin: 0; page-break-after: always; }
        }
    </style>
</head>
<body>
    <div id="header-bar">
        <h1>📊 PowerPoint Preview ({} slides)</h1>
        <div class="header-actions">
            <button class="header-btn" onclick="window.print()">🖨️ Print</button>
            <button class="header-btn" onclick="window.close()">✖️ Close</button>
        </div>
    </div>
    <div class="container">
""".format(slide_count))
        
        # Process each slide
        for slide_num, slide in enumerate(prs.slides, start=1):
            html_parts.append(f'<div class="slide">')
            html_parts.append(f'<div class="slide-header">')
            html_parts.append(f'<div class="slide-number">Slide {slide_num} of {slide_count}</div>')
            
            # Try to get slide title
            title_text = ""
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                html_parts.append(f'<div class="slide-title">{_escape_html(title_text)}</div>')
            
            html_parts.append(f'</div>')
            html_parts.append(f'<div class="slide-content">')
            
            # Process all shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if text and text != title_text:  # Don't duplicate title
                        html_parts.append(f'<div class="text-box">')
                        # Handle paragraphs
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                level = paragraph.level
                                indent = level * 20
                                html_parts.append(f'<p style="margin-left: {indent}px;">{_escape_html(paragraph.text)}</p>')
                        html_parts.append(f'</div>')
            
            html_parts.append(f'</div>')
            html_parts.append(f'</div>')
        
        html_parts.append("""
    </div>
</body>
</html>
""")
        
        # Write HTML file
        html_content = ''.join(html_parts)
        with open(html_file, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        print(f"HTML file created successfully: {len(html_content)} bytes")
        return True
        
    except Exception as e:
        print(f"ERROR: python-pptx conversion failed: {e}")
        traceback.print_exc()
        return False

def _escape_html(text):
    """Escape HTML special characters."""
    if not text:
        return ""
    return (text
        .replace('&', '&amp;')
        .replace('<', '&lt;')
        .replace('>', '&gt;')
        .replace('"', '&quot;')
        .replace("'", '&#39;'))

def convert_ppt_to_html_libreoffice(ppt_file, html_file):
    """
    Convert PPT/PPTX to HTML using LibreOffice.
    
    Args:
        ppt_file (str): Path to input PPT/PPTX file
        html_file (str): Path to output HTML file
    
    Returns:
        bool: True if conversion successful, False otherwise
    """
    print(f"Attempting PPT/PPTX to HTML conversion with LibreOffice...")
    
    try:
        # Ensure output directory exists
        output_dir = os.path.dirname(html_file)
        os.makedirs(output_dir, exist_ok=True)
        
        # Verify input file exists and check its properties
        if not os.path.exists(ppt_file):
            raise FileNotFoundError(f"Input file not found: {ppt_file}")
        
        file_size = os.path.getsize(ppt_file)
        print(f"Input file size: {file_size} bytes ({file_size / 1024 / 1024:.2f} MB)")
        
        # Read first few bytes to verify file signature
        with open(ppt_file, 'rb') as f:
            magic = f.read(8)
            print(f"File signature (hex): {magic.hex()}")
        
        # Detect file type from extension and signature
        file_ext = os.path.splitext(ppt_file)[1].lower()
        print(f"File extension: {file_ext}")
        
        # PPTX files should start with PK (ZIP signature: 50 4B)
        # PPT files should start with D0 CF 11 E0 (OLE2 signature)
        is_pptx = magic[:2] == b'PK'
        is_ppt = magic[:4] == b'\xD0\xCF\x11\xE0'
        
        print(f"File type detection: PPTX={is_pptx}, PPT={is_ppt}")
        
        # Set environment variables for LibreOffice
        env = os.environ.copy()
        env['SAL_USE_VCLPLUGIN'] = 'svp'
        env['HOME'] = os.path.expanduser('~')
        
        # Try multiple LibreOffice command variations
        # Method 1: Direct conversion with impress filter
        cmd = [
            'libreoffice',
            '--headless',
            '--invisible',
            '--nocrashreport',
            '--nodefault',
            '--nofirststartwizard',
            '--nolockcheck',
            '--nologo',
            '--norestore',
            '-env:UserInstallation=file:///tmp/libreoffice_user_profile',
            '--convert-to', 'html:impress_html_Export',
            '--outdir', output_dir,
            ppt_file
        ]
        
        print(f"Executing LibreOffice command: {' '.join(cmd)}")
        print(f"Environment: SAL_USE_VCLPLUGIN={env.get('SAL_USE_VCLPLUGIN')}, HOME={env.get('HOME')}")
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120, env=env)
        print("LibreOffice stdout:", result.stdout)
        print("LibreOffice stderr:", result.stderr)
        print("LibreOffice return code:", result.returncode)
        
        # If first method failed, try without filter specification
        if result.returncode != 0:
            print("First method failed, trying without filter...")
            cmd = [
                'libreoffice',
                '--headless',
                '--invisible',
                '--nocrashreport',
                '--nodefault',
                '--nofirststartwizard',
                '--nolockcheck',
                '--nologo',
                '--norestore',
                '-env:UserInstallation=file:///tmp/libreoffice_user_profile',
                '--convert-to', 'html',
                '--outdir', output_dir,
                ppt_file
            ]
            print(f"Executing LibreOffice command (fallback): {' '.join(cmd)}")
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120, env=env)
            print("LibreOffice stdout:", result.stdout)
            print("LibreOffice stderr:", result.stderr)
            print("LibreOffice return code:", result.returncode)
        
        # LibreOffice creates output with the input filename + .html extension
        base_name = os.path.splitext(os.path.basename(ppt_file))[0]
        html_output_file = os.path.join(output_dir, f"{base_name}.html")
        
        if not os.path.exists(html_output_file):
            # Try to find any .html file in the output directory
            html_files = [f for f in os.listdir(output_dir) if f.endswith('.html')]
            if html_files:
                html_output_file = os.path.join(output_dir, html_files[0])
            else:
                print(f"ERROR: No HTML file found in {output_dir}")
                print(f"Directory contents: {os.listdir(output_dir)}")
                raise FileNotFoundError(f"No HTML file produced by LibreOffice in {output_dir}")

        # Rename to desired output filename if different
        if html_output_file != html_file:
            shutil.move(html_output_file, html_file)
        
        file_size = os.path.getsize(html_file)
        print(f"HTML file created successfully: {file_size} bytes")
        
        # Enhance HTML with header bar for print/close
        enhance_html_with_header(html_file)
        
        return True
        
    except FileNotFoundError as e:
        print(f"ERROR: LibreOffice not found or file error: {e}")
        return False
    except subprocess.TimeoutExpired:
        print("ERROR: LibreOffice conversion timed out (>120s)")
        return False
    except Exception as e:
        print(f"ERROR: LibreOffice conversion error: {e}")
        traceback.print_exc()
        return False

def enhance_html_with_header(html_file):
    """Add header bar with print and close buttons to HTML."""
    try:
        with open(html_file, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        header_html = """
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PowerPoint Preview</title>
    <style>
        body { margin: 0; padding: 0; font-family: Arial, sans-serif; }
        #header-bar {
            position: sticky;
            top: 0;
            background: linear-gradient(135deg, #0078d4 0%, #00a4ef 100%);
            color: white;
            padding: 12px 20px;
            display: flex;
            justify-content: space-between;
            align-items: center;
            box-shadow: 0 2px 8px rgba(0,0,0,0.1);
            z-index: 1000;
        }
        #header-bar h1 {
            margin: 0;
            font-size: 18px;
            font-weight: 600;
        }
        .header-actions {
            display: flex;
            gap: 10px;
        }
        .header-btn {
            background: rgba(255,255,255,0.2);
            border: 1px solid rgba(255,255,255,0.3);
            color: white;
            padding: 8px 16px;
            border-radius: 6px;
            cursor: pointer;
            font-size: 14px;
            font-weight: 500;
            transition: all 0.2s;
        }
        .header-btn:hover {
            background: rgba(255,255,255,0.3);
            transform: translateY(-1px);
        }
        #content-area {
            padding: 20px;
            max-width: 1200px;
            margin: 0 auto;
        }
        @media print {
            #header-bar { display: none; }
            #content-area { padding: 0; max-width: none; }
        }
    </style>
</head>
<body>
    <div id="header-bar">
        <h1>📊 PowerPoint Preview</h1>
        <div class="header-actions">
            <button class="header-btn" onclick="window.print()">🖨️ Print</button>
            <button class="header-btn" onclick="window.close()">✖️ Close</button>
        </div>
    </div>
    <div id="content-area">
"""
        
        footer_html = """
    </div>
</body>
</html>
"""
        
        # Extract body content if it exists
        if '<body' in content.lower():
            body_start = content.lower().find('<body')
            body_end = content.lower().find('</body>')
            if body_start != -1 and body_end != -1:
                body_content = content[content.find('>', body_start) + 1:body_end]
                enhanced_content = header_html + body_content + footer_html
            else:
                enhanced_content = header_html + content + footer_html
        else:
            enhanced_content = header_html + content + footer_html
        
        with open(html_file, 'w', encoding='utf-8') as f:
            f.write(enhanced_content)
        
        print("HTML enhanced with header bar")
    except Exception as e:
        print(f"Warning: Could not enhance HTML: {e}")

def convert_ppt_to_html_unoconv(ppt_file, html_file):
    """
    Convert PPT/PPTX to HTML using unoconv (alternative method).
    
    Args:
        ppt_file (str): Path to input PPT/PPTX file
        html_file (str): Path to output HTML file
    
    Returns:
        bool: True if conversion successful, False otherwise
    """
    print(f"Attempting PPT/PPTX to HTML conversion with unoconv...")
    
    try:
        output_dir = os.path.dirname(html_file)
        os.makedirs(output_dir, exist_ok=True)
        
        cmd = [
            'unoconv',
            '-f', 'html',
            '-o', html_file,
            ppt_file
        ]
        
        print(f"Executing unoconv command: {' '.join(cmd)}")
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        print("unoconv stdout:", result.stdout)
        print("unoconv stderr:", result.stderr)
        print("unoconv return code:", result.returncode)
        
        if result.returncode == 0 and os.path.exists(html_file):
            print(f"HTML file created successfully: {os.path.getsize(html_file)} bytes")
            enhance_html_with_header(html_file)
            return True
        else:
            print("ERROR: unoconv conversion failed")
            return False
            
    except FileNotFoundError:
        print("ERROR: unoconv not found")
        return False
    except Exception as e:
        print(f"ERROR: unoconv conversion error: {e}")
        traceback.print_exc()
        return False

def create_error_html(html_file, error_message):
    """Create a simple error HTML file when conversion fails."""
    try:
        html_content = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>PowerPoint Preview Error</title>
    <style>
        body {{
            font-family: Arial, sans-serif;
            max-width: 800px;
            margin: 50px auto;
            padding: 20px;
            background: #f5f5f5;
        }}
        .error-box {{
            background: white;
            border-left: 4px solid #d32f2f;
            padding: 20px;
            border-radius: 4px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        }}
        h1 {{
            color: #d32f2f;
            margin-top: 0;
        }}
        .error-message {{
            background: #ffebee;
            padding: 15px;
            border-radius: 4px;
            margin: 15px 0;
            font-family: monospace;
            font-size: 14px;
        }}
        .suggestion {{
            background: #e3f2fd;
            padding: 15px;
            border-radius: 4px;
            border-left: 4px solid #2196f3;
        }}
    </style>
</head>
<body>
    <div class="error-box">
        <h1>⚠️ PowerPoint Preview Error</h1>
        <p>The PowerPoint file could not be converted to HTML for preview.</p>
        
        <div class="error-message">
            {error_message}
        </div>
        
        <div class="suggestion">
            <h3>💡 Suggestions:</h3>
            <ul>
                <li>Try downloading the file directly using the download button</li>
                <li>Ensure the file is a valid PowerPoint presentation (.ppt or .pptx)</li>
                <li>The file may contain features not supported by the preview converter</li>
                <li>Very large presentations may timeout during conversion</li>
            </ul>
        </div>
    </div>
</body>
</html>
"""
        with open(html_file, 'w', encoding='utf-8') as f:
            f.write(html_content)
        print(f"Created error HTML file: {html_file}")
        return True
    except Exception as e:
        print(f"ERROR: Could not create error HTML: {e}")
        return False

def main():
    parser = argparse.ArgumentParser(description='Convert PPT/PPTX to HTML for web preview')
    parser.add_argument('ppt_file', help='Input PPT/PPTX file path')
    parser.add_argument('html_file', help='Output HTML file path')
    
    args = parser.parse_args()
    
    print("=== PPT/PPTX to HTML Converter ===")
    print(f"Python version: {sys.version}")
    print(f"Working directory: {os.getcwd()}")
    print(f"Arguments: {vars(args)}")
    
    # Check if input file exists
    if not os.path.exists(args.ppt_file):
        print(f"ERROR: Input PPT/PPTX file not found: {args.ppt_file}")
        sys.exit(1)
    
    # Determine file type
    file_ext = os.path.splitext(args.ppt_file)[1].lower()
    if file_ext == '.ppt':
        print("Detected format: PPT (PowerPoint 97-2003)")
    elif file_ext == '.pptx':
        print("Detected format: PPTX (PowerPoint 2007+)")
    else:
        print(f"WARNING: Unknown extension '{file_ext}', will try to convert anyway")
    
    # Create output directory if it doesn't exist
    output_dir = os.path.dirname(args.html_file)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    success = False
    
    # For PPTX files, try python-pptx first (fastest and most reliable)
    if file_ext == '.pptx' and PPTX_AVAILABLE:
        print("\n=== Trying Method 1: python-pptx (Direct Parsing) ===")
        success = convert_pptx_to_html_direct(args.ppt_file, args.html_file)
    
    # If python-pptx failed or not available, try LibreOffice
    if not success:
        print("\n=== Trying Method 2: LibreOffice ===")
        success = convert_ppt_to_html_libreoffice(args.ppt_file, args.html_file)
    
    # If LibreOffice failed, try unoconv
    if not success:
        print("\n=== Trying Method 3: unoconv ===")
        success = convert_ppt_to_html_unoconv(args.ppt_file, args.html_file)
    
    if success:
        print("=== CONVERSION SUCCESSFUL ===")
        sys.exit(0)
    else:
        print("=== CONVERSION FAILED ===")
        error_msg = "LibreOffice conversion failed with 'Unsupported document type' error. This usually means LibreOffice cannot process this file format or the file may be corrupted."
        
        # Create an error HTML file so at least something is returned
        if create_error_html(args.html_file, error_msg):
            print("Created error HTML page for user feedback")
            sys.exit(0)  # Exit with success so the error page is shown
        else:
            print("ERROR: All conversion methods failed. Please ensure LibreOffice or unoconv is installed.")
            print("Hint: The 'Unsupported document type' error often means LibreOffice can't access the file.")
            print("      Check file permissions and ensure the file is a valid PowerPoint document.")
            sys.exit(1)

if __name__ == "__main__":
    main()

