#!/usr/bin/env python3
"""
Optimized CSV to PPT Converter
Converts CSV files to PowerPoint (PPTX) format using pandas and python-pptx with performance optimizations.
"""

import pandas as pd
import argparse
import os
import sys
from datetime import datetime
import traceback

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError as e:
    print(f"ERROR: Required python-pptx library not available: {e}")
    print("Please install python-pptx: pip install python-pptx")
    sys.exit(1)

def create_ppt_from_csv_optimized(csv_file, output_file, title="CSV Data", author="Unknown", max_rows_per_slide=100):
    """
    Convert CSV file to PPTX format with performance optimizations.
    
    Args:
        csv_file (str): Path to input CSV file
        output_file (str): Path to output PPTX file
        title (str): Document title
        author (str): Document author
        max_rows_per_slide (int): Maximum rows per slide (increased for better performance)
    """
    print(f"Starting optimized CSV to PPT conversion...")
    print(f"Input: {csv_file}")
    print(f"Output: {output_file}")
    print(f"Title: {title}")
    print(f"Author: {author}")
    print(f"Max rows per slide: {max_rows_per_slide}")
    
    try:
        # Read CSV file with optimizations
        print("Reading CSV file...")
        df = pd.read_csv(csv_file, dtype=str)  # Read all as strings for speed
        print(f"CSV loaded successfully: {len(df)} rows, {len(df.columns)} columns")
        
        # Limit processing for very large files
        max_total_rows = 10000  # Limit to 10k rows for performance
        if len(df) > max_total_rows:
            print(f"Large dataset detected ({len(df)} rows). Limiting to first {max_total_rows} rows for performance.")
            df = df.head(max_total_rows)
        
        # Create PowerPoint presentation
        print("Creating PowerPoint presentation...")
        prs = Presentation()
        
        # Set presentation properties
        prs.core_properties.title = title
        prs.core_properties.author = author
        prs.core_properties.created = datetime.now()
        
        # Add title slide
        print("Adding title slide...")
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = f"Data Analysis Report\nGenerated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\nAuthor: {author}"
        
        # Add summary slide
        print("Adding summary slide...")
        content_slide_layout = prs.slide_layouts[1]
        summary_slide = prs.slides.add_slide(content_slide_layout)
        
        summary_title = summary_slide.shapes.title
        summary_title.text = "Data Summary"
        
        summary_content = summary_slide.placeholders[1]
        summary_text = f"• Total Rows: {len(df):,}\n• Total Columns: {len(df.columns)}\n• File Size: {os.path.getsize(csv_file):,} bytes\n• Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        summary_content.text = summary_text
        
        # Add column headers slide
        print("Adding column headers slide...")
        headers_slide = prs.slides.add_slide(content_slide_layout)
        headers_title = headers_slide.shapes.title
        headers_title.text = "Column Headers"
        
        headers_content = headers_slide.placeholders[1]
        headers_text = "\n".join([f"• {i+1}. {col}" for i, col in enumerate(df.columns)])
        headers_content.text = headers_text
        
        # Process data in larger chunks for better performance
        chunk_size = max_rows_per_slide
        total_chunks = (len(df) + chunk_size - 1) // chunk_size
        
        print(f"Creating {total_chunks} data slides...")
        
        # Pre-calculate table dimensions
        cols = len(df.columns)
        table_width = Inches(9)
        col_width = table_width / cols
        
        for chunk_idx in range(total_chunks):
            start_idx = chunk_idx * chunk_size
            end_idx = min((chunk_idx + 1) * chunk_size, len(df))
            chunk_df = df.iloc[start_idx:end_idx]
            
            if chunk_idx % 10 == 0:  # Progress update every 10 slides
                print(f"Processing chunk {chunk_idx + 1}/{total_chunks} (rows {start_idx + 1}-{end_idx})")
            
            # Create data slide
            data_slide = prs.slides.add_slide(content_slide_layout)
            data_title = data_slide.shapes.title
            data_title.text = f"Data - Part {chunk_idx + 1} of {total_chunks}"
            
            # Create table with optimized dimensions
            rows = len(chunk_df) + 1  # +1 for header row
            
            # Calculate optimal table size
            table_height = min(Inches(5), Inches(0.3) * rows)  # Dynamic height based on rows
            
            table = data_slide.shapes.add_table(
                rows, cols, 
                Inches(0.5), Inches(1.5), 
                table_width, table_height
            ).table
            
            # Set column widths
            for col_idx in range(cols):
                table.columns[col_idx].width = col_width
            
            # Add header row with optimized styling
            for col_idx, col_name in enumerate(chunk_df.columns):
                cell = table.cell(0, col_idx)
                cell.text = str(col_name)
                
                # Optimized header styling
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = True
                paragraph.font.size = Pt(9)  # Smaller font for more data
                paragraph.alignment = PP_ALIGN.CENTER
                
                # Set background color
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                
                # Set text color
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
            
            # Add data rows with optimized processing
            for row_idx, (_, row) in enumerate(chunk_df.iterrows(), 1):
                for col_idx, value in enumerate(row):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(value) if pd.notna(value) else ""
                    
                    # Optimized data cell styling
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(8)  # Smaller font for more data
                    paragraph.alignment = PP_ALIGN.LEFT
                    
                    # Alternating row colors
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
        
        # Add conclusion slide
        print("Adding conclusion slide...")
        conclusion_slide = prs.slides.add_slide(content_slide_layout)
        conclusion_title = conclusion_slide.shapes.title
        conclusion_title.text = "Conclusion"
        
        conclusion_content = conclusion_slide.placeholders[1]
        conclusion_text = f"Data processing complete.\n\n• Total rows processed: {len(df):,}\n• Total columns: {len(df.columns)}\n• Slides created: {total_chunks + 3}\n• Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        conclusion_content.text = conclusion_text
        
        # Save presentation
        print("Saving presentation...")
        prs.save(output_file)
        
        # Verify file was created
        if os.path.exists(output_file):
            file_size = os.path.getsize(output_file)
            print(f"PPT file created successfully: {file_size} bytes")
            print(f"Total slides: {total_chunks + 3}")
            return True
        else:
            print("ERROR: PPT file was not created")
            return False
            
    except Exception as e:
        print(f"ERROR: Failed to create PPT from CSV: {e}")
        traceback.print_exc()
        return False

def main():
    parser = argparse.ArgumentParser(description='Convert CSV to PPT format (optimized)')
    parser.add_argument('csv_file', help='Input CSV file path')
    parser.add_argument('output_file', help='Output PPT file path')
    parser.add_argument('--title', default='CSV Data', help='Document title')
    parser.add_argument('--author', default='Unknown', help='Document author')
    parser.add_argument('--max-rows', type=int, default=100, help='Maximum rows per slide')
    
    args = parser.parse_args()
    
    print("=== Optimized CSV to PPT Converter ===")
    print(f"Python version: {sys.version}")
    print(f"Working directory: {os.getcwd()}")
    print(f"Arguments: {vars(args)}")
    
    # Check if input file exists
    if not os.path.exists(args.csv_file):
        print(f"ERROR: Input CSV file not found: {args.csv_file}")
        sys.exit(1)
    
    # Check pandas availability
    try:
        print(f"Pandas version: {pd.__version__}")
    except Exception as e:
        print(f"ERROR: Pandas not available: {e}")
        sys.exit(1)
    
    # Create output directory if it doesn't exist
    output_dir = os.path.dirname(args.output_file)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    # Convert CSV to PPT
    success = create_ppt_from_csv_optimized(
        args.csv_file,
        args.output_file,
        args.title,
        args.author,
        args.max_rows
    )
    
    if success:
        print("=== CONVERSION SUCCESSFUL ===")
        sys.exit(0)
    else:
        print("=== CONVERSION FAILED ===")
        sys.exit(1)

if __name__ == "__main__":
    main()
