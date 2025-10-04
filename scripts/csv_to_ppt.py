#!/usr/bin/env python3
"""
CSV to PPT Converter
Converts CSV files to PowerPoint (PPTX) format using pandas and python-pptx.
"""

import pandas as pd
import argparse
import os
import sys
from datetime import datetime
import traceback

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError as e:
    print(f"ERROR: Required python-pptx library not available: {e}")
    print("Please install python-pptx: pip install python-pptx")
    sys.exit(1)

def create_ppt_from_csv(csv_file, output_file, title="CSV Data", author="Unknown", max_rows_per_slide=50):
    """
    Convert CSV file to PPTX format.
    
    Args:
        csv_file (str): Path to input CSV file
        output_file (str): Path to output PPTX file
        title (str): Document title
        author (str): Document author
        max_rows_per_slide (int): Maximum rows per slide to prevent overcrowding
    """
    print(f"Starting CSV to PPT conversion...")
    print(f"Input: {csv_file}")
    print(f"Output: {output_file}")
    print(f"Title: {title}")
    print(f"Author: {author}")
    print(f"Max rows per slide: {max_rows_per_slide}")
    
    try:
        # Read CSV file
        print("Reading CSV file...")
        df = pd.read_csv(csv_file)
        print(f"CSV loaded successfully: {len(df)} rows, {len(df.columns)} columns")
        
        # Process all rows - no limits
        print(f"Processing all {len(df)} rows (including any repeated data)")
        
        # Create PowerPoint presentation
        print("Creating PowerPoint presentation...")
        prs = Presentation()
        
        # Set presentation properties
        prs.core_properties.title = title
        prs.core_properties.author = author
        prs.core_properties.created = datetime.now()
        
        # Add title slide
        print("Adding title slide...")
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = f"Data Analysis Report\nGenerated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\nAuthor: {author}"
        
        # Add summary slide
        print("Adding summary slide...")
        content_slide_layout = prs.slide_layouts[1]  # Title and content layout
        summary_slide = prs.slides.add_slide(content_slide_layout)
        
        summary_title = summary_slide.shapes.title
        summary_title.text = "Data Summary"
        
        summary_content = summary_slide.placeholders[1]
        summary_text = f"• Total Rows: {len(df):,}\n• Total Columns: {len(df.columns)}\n• File Size: {os.path.getsize(csv_file):,} bytes\n• Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        summary_content.text = summary_text
        
        # Process data in chunks for slides
        print("Processing data for slides...")
        total_slides = 0
        
        # Add column headers slide
        print("Adding column headers slide...")
        headers_slide = prs.slides.add_slide(content_slide_layout)
        headers_title = headers_slide.shapes.title
        headers_title.text = "Column Headers"
        
        headers_content = headers_slide.placeholders[1]
        headers_text = "\n".join([f"• {i+1}. {col}" for i, col in enumerate(df.columns)])
        headers_content.text = headers_text
        total_slides += 1
        
        # Process data in chunks
        chunk_size = max_rows_per_slide
        total_chunks = (len(df) + chunk_size - 1) // chunk_size
        
        print(f"Creating {total_chunks} data slides...")
        
        for chunk_idx in range(total_chunks):
            start_idx = chunk_idx * chunk_size
            end_idx = min((chunk_idx + 1) * chunk_size, len(df))
            chunk_df = df.iloc[start_idx:end_idx]
            
            print(f"Processing chunk {chunk_idx + 1}/{total_chunks} (rows {start_idx + 1}-{end_idx})")
            
            # Create data slide
            data_slide = prs.slides.add_slide(content_slide_layout)
            data_title = data_slide.shapes.title
            data_title.text = f"Data - Part {chunk_idx + 1} of {total_chunks}"
            
            # Create table
            rows, cols = len(chunk_df) + 1, len(chunk_df.columns)  # +1 for header row
            
            # Calculate table position and size
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            
            table = data_slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Set column widths
            for col_idx in range(cols):
                table.columns[col_idx].width = Inches(9 / cols)
            
            # Add header row
            for col_idx, col_name in enumerate(chunk_df.columns):
                cell = table.cell(0, col_idx)
                cell.text = str(col_name)
                
                # Style header cell
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196)  # Blue header
                
                # Set text color to white
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)
            
            # Add data rows
            for row_idx, (_, row) in enumerate(chunk_df.iterrows(), 1):
                for col_idx, value in enumerate(row):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(value) if pd.notna(value) else ""
                    
                    # Style data cell
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                    
                    # Alternate row colors
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(242, 242, 242)  # Light gray
            
            total_slides += 1
        
        # Add conclusion slide
        print("Adding conclusion slide...")
        conclusion_slide = prs.slides.add_slide(content_slide_layout)
        conclusion_title = conclusion_slide.shapes.title
        conclusion_title.text = "Conclusion"
        
        conclusion_content = conclusion_slide.placeholders[1]
        conclusion_text = f"Data Processing Complete\n\n• Total slides created: {total_slides + 2}\n• Data rows processed: {len(df):,}\n• Columns analyzed: {len(df.columns)}\n• File generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        conclusion_content.text = conclusion_text
        
        # Save presentation
        print(f"Saving PowerPoint presentation to {output_file}...")
        prs.save(output_file)
        
        # Verify file was created
        if os.path.exists(output_file):
            file_size = os.path.getsize(output_file)
            print(f"PPTX file created successfully: {file_size} bytes")
            print(f"Total slides created: {total_slides + 2}")
            return True
        else:
            print("ERROR: PPTX file was not created")
            return False
            
    except Exception as e:
        print(f"ERROR: Failed to create PPTX from CSV: {e}")
        traceback.print_exc()
        return False

def main():
    parser = argparse.ArgumentParser(description='Convert CSV to PPTX format')
    parser.add_argument('csv_file', help='Input CSV file path')
    parser.add_argument('output_file', help='Output PPTX file path')
    parser.add_argument('--title', default='CSV Data', help='Document title')
    parser.add_argument('--author', default='Unknown', help='Document author')
    parser.add_argument('--max-rows-per-slide', type=int, default=50, help='Maximum rows per slide')
    
    args = parser.parse_args()
    
    print("=== CSV to PPTX Converter ===")
    print(f"Python version: {sys.version}")
    print(f"Working directory: {os.getcwd()}")
    print(f"Arguments: {vars(args)}")
    
    # Check if input file exists
    if not os.path.exists(args.csv_file):
        print(f"ERROR: Input CSV file not found: {args.csv_file}")
        sys.exit(1)
    
    # Check pandas availability
    try:
        print(f"Pandas version: {pd.__version__}")
    except Exception as e:
        print(f"ERROR: Pandas not available: {e}")
        sys.exit(1)
    
    # Check python-pptx availability
    try:
        from pptx import __version__ as pptx_version
        print(f"python-pptx version: {pptx_version}")
    except Exception as e:
        print(f"ERROR: python-pptx not available: {e}")
        sys.exit(1)
    
    # Create output directory if it doesn't exist
    output_dir = os.path.dirname(args.output_file)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    # Convert CSV to PPTX
    success = create_ppt_from_csv(
        args.csv_file,
        args.output_file,
        args.title,
        args.author,
        args.max_rows_per_slide
    )
    
    if success:
        print("=== CONVERSION SUCCESSFUL ===")
        sys.exit(0)
    else:
        print("=== CONVERSION FAILED ===")
        sys.exit(1)

if __name__ == "__main__":
    main()
